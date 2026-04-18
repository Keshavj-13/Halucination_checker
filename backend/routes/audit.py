from fastapi import APIRouter, HTTPException, UploadFile, File
from models.schemas import AuditRequest, AuditResponse
from services.claim_extractor import extract_claims
from services.verifier import verify_claims, verify_claims_stream
from data.sample_data import SAMPLE_RESPONSE
import logging
from fastapi.responses import StreamingResponse
import json
import io
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

router = APIRouter(prefix="/audit", tags=["audit"])
logger = logging.getLogger("audit-api.routes")

def extract_text_from_file(file: UploadFile) -> str:
    """Extract text content from various file formats."""
    filename = file.filename.lower()
    logger.info(f"Starting text extraction for file: {file.filename}")
    
    # Read file content
    content = file.file.read()
    logger.info(f"File read successfully, size: {len(content)} bytes")
    
    if filename.endswith('.pdf'):
        logger.info("Detected PDF file, extracting text from pages")
        # PDF extraction
        pdf_reader = PdfReader(io.BytesIO(content))
        text = ""
        for i, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            text += page_text + "\n"
            logger.debug(f"Extracted text from page {i+1}: {len(page_text)} characters")
        logger.info(f"PDF text extraction completed, total text length: {len(text)}")
        return text
    
    elif filename.endswith('.docx'):
        logger.info("Detected DOCX file, extracting text from paragraphs")
        # DOCX extraction
        doc = Document(io.BytesIO(content))
        text = ""
        for i, paragraph in enumerate(doc.paragraphs):
            para_text = paragraph.text
            text += para_text + "\n"
            logger.debug(f"Extracted text from paragraph {i+1}: {len(para_text)} characters")
        logger.info(f"DOCX text extraction completed, total text length: {len(text)}")
        return text
    
    elif filename.endswith('.pptx'):
        logger.info("Detected PPTX file, extracting text from slides")
        # PPTX extraction
        prs = Presentation(io.BytesIO(content))
        text = ""
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            text += slide_text
            logger.debug(f"Extracted text from slide {i+1}: {len(slide_text)} characters")
        logger.info(f"PPTX text extraction completed, total text length: {len(text)}")
        return text
    
    else:
        logger.info("Detected plain text file, attempting to decode")
        # Try to decode as text (for .txt, .md, .html, etc.)
        try:
            text = content.decode('utf-8')
            logger.info(f"UTF-8 text decoding successful, text length: {len(text)}")
            return text
        except UnicodeDecodeError:
            try:
                text = content.decode('latin-1')  # Fallback encoding
                logger.info(f"Latin-1 text decoding successful, text length: {len(text)}")
                return text
            except:
                logger.error(f"Failed to decode text from file: {file.filename}")
                raise HTTPException(status_code=400, detail="Unsupported file format or encoding")

@router.post("/", response_model=AuditResponse)
async def run_audit(body: AuditRequest):
    logger.info(f"Received audit request for document length: {len(body.document)}")
    
    # 1. Extract Claims (with indices)
    claims_data = await extract_claims(body.document)
    
    # 2. Verify Claims (Multilayer)
    verified_claims = await verify_claims(claims_data)
    
    # 3. Calculate Stats
    verified = len([c for c in verified_claims if c.status == "Verified"])
    plausible = len([c for c in verified_claims if c.status == "Plausible"])
    hallucinations = len([c for c in verified_claims if c.status == "Hallucination"])
    
    return AuditResponse(
        document=body.document,
        total=len(verified_claims),
        verified=verified,
        plausible=plausible,
        hallucinations=hallucinations,
        claims=verified_claims
    )

@router.post("/stream")
async def run_audit_stream(body: AuditRequest):
    logger.info(f"Received streaming audit request")
    claims_data = await extract_claims(body.document)
    
    async def event_generator():
        yield f"data: {json.dumps({'type': 'start', 'total': len(claims_data)})}\n\n"
        
        async for claim in verify_claims_stream(claims_data):
            yield f"data: {json.dumps({'type': 'claim', 'claim': claim.model_dump()})}\n\n"
            
        yield f"data: {json.dumps({'type': 'done'})}\n\n"

    return StreamingResponse(event_generator(), media_type="text/event-stream")

@router.post("/upload", response_model=AuditResponse)
async def upload_and_audit(file: UploadFile = File(...)):
    logger.info(f"Received file upload: {file.filename}")
    
    # Extract text from file
    document_text = extract_text_from_file(file)
    
    if not document_text.strip():
        logger.warning(f"No text content found in file: {file.filename}")
        raise HTTPException(status_code=400, detail="No text content found in file")
    
    logger.info(f"Extracted text length: {len(document_text)}")
    
    # Process like regular audit
    logger.info("Starting claim extraction")
    claims_data = await extract_claims(document_text)
    logger.info(f"Extracted {len(claims_data)} claims")
    
    logger.info("Starting claim verification")
    verified_claims = await verify_claims(claims_data)
    logger.info(f"Verified {len(verified_claims)} claims")
    
    verified = len([c for c in verified_claims if c.status == "Verified"])
    plausible = len([c for c in verified_claims if c.status == "Plausible"])
    hallucinations = len([c for c in verified_claims if c.status == "Hallucination"])
    
    logger.info(f"Audit completed - Verified: {verified}, Plausible: {plausible}, Hallucinations: {hallucinations}")
    
    return AuditResponse(
        document=document_text,
        total=len(verified_claims),
        verified=verified,
        plausible=plausible,
        hallucinations=hallucinations,
        claims=verified_claims
    )

@router.post("/upload/stream")
async def upload_and_audit_stream(file: UploadFile = File(...)):
    logger.info(f"Received streaming file upload: {file.filename}")
    
    # Extract text from file
    document_text = extract_text_from_file(file)
    
    if not document_text.strip():
        raise HTTPException(status_code=400, detail="No text content found in file")
    
    claims_data = await extract_claims(document_text)
    
    async def event_generator():
        yield f"data: {json.dumps({'type': 'start', 'total': len(claims_data)})}\n\n"
        
        async for claim in verify_claims_stream(claims_data):
            yield f"data: {json.dumps({'type': 'claim', 'claim': claim.model_dump()})}\n\n"
            
        yield f"data: {json.dumps({'type': 'done'})}\n\n"

    return StreamingResponse(event_generator(), media_type="text/event-stream")

@router.get("/sample", response_model=AuditResponse)
def get_sample():
    return SAMPLE_RESPONSE
